#!/usr/bin/env python3
"""
Notion Hardware Page → PPTX Presentation Generator

Usage:
    pip install python-pptx requests
    python generate_presentation.py [--template template.pptx] [--output output.pptx]

Environment / Defaults:
    NOTION_API_KEY: Notion integration token (default provided)
    NOTION_PAGE_ID: Notion page ID (default provided)
"""

import argparse
import json
import os
import sys
from pathlib import Path

import requests
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ── Notion API ──────────────────────────────────────────────────────────────

NOTION_API_KEY = os.getenv("NOTION_API_KEY", "")
NOTION_PAGE_ID = os.getenv("NOTION_PAGE_ID", "31a607db92b680ada380efe47a09712b")
NOTION_API_URL = "https://api.notion.com/v1"
NOTION_HEADERS = {
    "Authorization": f"Bearer {NOTION_API_KEY}",
    "Notion-Version": "2022-06-28",
    "Content-Type": "application/json",
}


def notion_get(endpoint: str, params: dict | None = None) -> dict:
    url = f"{NOTION_API_URL}/{endpoint}"
    resp = requests.get(url, headers=NOTION_HEADERS, params=params, timeout=30)
    resp.raise_for_status()
    return resp.json()


def get_page_title(page_id: str) -> str:
    data = notion_get(f"pages/{page_id}")
    props = data.get("properties", {})
    for prop in props.values():
        if prop.get("type") == "title":
            titles = prop.get("title", [])
            return "".join(t.get("plain_text", "") for t in titles)
    return "Hardware Presentation"


def get_block_children(block_id: str) -> list[dict]:
    blocks: list[dict] = []
    cursor = None
    while True:
        params = {"page_size": 100}
        if cursor:
            params["start_cursor"] = cursor
        data = notion_get(f"blocks/{block_id}/children", params)
        blocks.extend(data.get("results", []))
        if not data.get("has_more"):
            break
        cursor = data.get("next_cursor")
    return blocks


def extract_rich_text(rich_texts: list[dict]) -> str:
    return "".join(rt.get("plain_text", "") for rt in rich_texts)


def parse_table_block(block: dict) -> list[list[str]]:
    """Fetch table rows from a table block."""
    rows_data = get_block_children(block["id"])
    table: list[list[str]] = []
    for row_block in rows_data:
        if row_block["type"] != "table_row":
            continue
        cells = row_block["table_row"]["cells"]
        table.append([extract_rich_text(cell) for cell in cells])
    return table


# ── Content Parsing ─────────────────────────────────────────────────────────

def parse_notion_page(page_id: str) -> list[dict]:
    """
    Parse the Notion page into a list of sections.
    Each section = { "title": str, "content": list[str | dict] }
    A dict item means a table: {"type": "table", "headers": [...], "rows": [[...]]}
    """
    blocks = get_block_children(page_id)
    sections: list[dict] = []
    current_section: dict | None = None

    for block in blocks:
        btype = block["type"]

        # Heading → new section
        if btype in ("heading_1", "heading_2", "heading_3"):
            text = extract_rich_text(block[btype].get("rich_text", []))
            current_section = {"title": text, "content": []}
            sections.append(current_section)
            continue

        if current_section is None:
            current_section = {"title": "", "content": []}
            sections.append(current_section)

        if btype == "paragraph":
            text = extract_rich_text(block["paragraph"].get("rich_text", []))
            if text.strip():
                current_section["content"].append(text)

        elif btype == "bulleted_list_item":
            text = extract_rich_text(block["bulleted_list_item"].get("rich_text", []))
            if text.strip():
                current_section["content"].append(f"• {text}")

        elif btype == "numbered_list_item":
            text = extract_rich_text(block["numbered_list_item"].get("rich_text", []))
            if text.strip():
                current_section["content"].append(text)

        elif btype == "to_do":
            text = extract_rich_text(block["to_do"].get("rich_text", []))
            checked = block["to_do"].get("checked", False)
            mark = "☑" if checked else "☐"
            if text.strip():
                current_section["content"].append(f"{mark} {text}")

        elif btype == "toggle":
            text = extract_rich_text(block["toggle"].get("rich_text", []))
            if text.strip():
                current_section["content"].append(f"▸ {text}")
            # Fetch toggle children
            if block.get("has_children"):
                children = get_block_children(block["id"])
                for child in children:
                    ctype = child["type"]
                    if ctype == "paragraph":
                        ct = extract_rich_text(child["paragraph"].get("rich_text", []))
                        if ct.strip():
                            current_section["content"].append(f"  {ct}")
                    elif ctype == "bulleted_list_item":
                        ct = extract_rich_text(child["bulleted_list_item"].get("rich_text", []))
                        if ct.strip():
                            current_section["content"].append(f"  • {ct}")
                    elif ctype == "table":
                        table_data = parse_table_block(child)
                        if table_data:
                            current_section["content"].append({
                                "type": "table",
                                "headers": table_data[0] if table_data else [],
                                "rows": table_data[1:] if len(table_data) > 1 else [],
                            })

        elif btype == "table":
            table_data = parse_table_block(block)
            if table_data:
                current_section["content"].append({
                    "type": "table",
                    "headers": table_data[0] if table_data else [],
                    "rows": table_data[1:] if len(table_data) > 1 else [],
                })

        elif btype == "callout":
            text = extract_rich_text(block["callout"].get("rich_text", []))
            icon = block["callout"].get("icon", {})
            emoji = icon.get("emoji", "💡") if icon.get("type") == "emoji" else "💡"
            if text.strip():
                current_section["content"].append(f"{emoji} {text}")

        elif btype == "code":
            text = extract_rich_text(block["code"].get("rich_text", []))
            if text.strip():
                current_section["content"].append(f"[Code] {text}")

        elif btype == "divider":
            pass  # skip

        elif btype == "image":
            img = block["image"]
            url = ""
            if img.get("type") == "file":
                url = img["file"].get("url", "")
            elif img.get("type") == "external":
                url = img["external"].get("url", "")
            caption = extract_rich_text(img.get("caption", []))
            current_section["content"].append(f"[Image: {caption or url}]")

    return sections


# ── PPTX Generation ─────────────────────────────────────────────────────────

# Theme colors
COLOR_TITLE_BG = RGBColor(0x1A, 0x1A, 0x2E)    # dark navy
COLOR_ACCENT = RGBColor(0x00, 0xB4, 0xD8)       # cyan accent
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
COLOR_DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
COLOR_TABLE_HEADER = RGBColor(0x00, 0x76, 0xA3)
COLOR_TABLE_ROW_ALT = RGBColor(0xF0, 0xF8, 0xFF)


def add_table_to_slide(slide, table_data: dict, left: float, top: float, width: float):
    """Add a formatted table to a slide."""
    headers = table_data["headers"]
    rows = table_data["rows"]
    if not headers:
        return

    num_rows = len(rows) + 1  # +1 for header
    num_cols = len(headers)

    tbl = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(left), Inches(top),
        Inches(width), Inches(0.4 * num_rows)
    ).table

    # Set column widths evenly
    col_width = Emu(int(Inches(width) / num_cols))
    for i in range(num_cols):
        tbl.columns[i].width = col_width

    # Header row
    for ci, header in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_TABLE_HEADER
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(10)
            paragraph.font.bold = True
            paragraph.font.color.rgb = COLOR_WHITE
            paragraph.alignment = PP_ALIGN.CENTER

    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = val
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_TABLE_ROW_ALT
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.font.color.rgb = COLOR_DARK_TEXT


def create_title_slide(prs: Presentation, title: str, subtitle: str = ""):
    """Create the title (cover) slide."""
    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_TITLE_BG

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    # Accent line
    from pptx.util import Emu
    line = slide.shapes.add_shape(
        1, Inches(3.5), Inches(4.0), Inches(3), Pt(3)  # rectangle as line
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()

    # Subtitle
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.3), Inches(8), Inches(0.8))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = COLOR_LIGHT_GRAY
        p2.alignment = PP_ALIGN.CENTER


def create_section_slide(prs: Presentation, section: dict):
    """Create a content slide for a section."""
    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

    # Title bar
    title_bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(1.0)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = COLOR_TITLE_BG
    title_bar.line.fill.background()

    # Accent stripe
    accent_stripe = slide.shapes.add_shape(
        1, Inches(0), Inches(1.0), Inches(10), Pt(4)
    )
    accent_stripe.fill.solid()
    accent_stripe.fill.fore_color.rgb = COLOR_ACCENT
    accent_stripe.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = section["title"]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.LEFT

    # Content area
    content_top = 1.4
    text_items = []
    tables = []

    for item in section["content"]:
        if isinstance(item, dict) and item.get("type") == "table":
            tables.append(item)
        else:
            text_items.append(str(item))

    # Add text content
    if text_items:
        content_text = "\n".join(text_items)
        max_height = 5.5 if not tables else 2.5
        txBox2 = slide.shapes.add_textbox(
            Inches(0.6), Inches(content_top),
            Inches(8.8), Inches(max_height)
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True

        for i, line in enumerate(text_items):
            if i == 0:
                p = tf2.paragraphs[0]
            else:
                p = tf2.add_paragraph()

            p.text = line
            p.font.size = Pt(14)
            p.font.color.rgb = COLOR_DARK_TEXT
            p.space_after = Pt(6)

            # Bullet items slightly indented
            if line.startswith("•") or line.startswith("  •"):
                p.font.size = Pt(13)
                p.level = 1

        content_top += min(len(text_items) * 0.35, max_height) + 0.2

    # Add tables
    for tbl_data in tables:
        add_table_to_slide(slide, tbl_data, 0.5, content_top, 9.0)
        num_rows = len(tbl_data.get("rows", [])) + 1
        content_top += num_rows * 0.4 + 0.3


def create_toc_slide(prs: Presentation, sections: list[dict]):
    """Create a table of contents slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_TITLE_BG

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "목차 / Contents"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    # Accent line
    line = slide.shapes.add_shape(
        1, Inches(0.5), Inches(1.2), Inches(2), Pt(3)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()

    # Section items
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(5.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True

    for i, section in enumerate(sections):
        if not section["title"]:
            continue
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = f"{i + 1}.  {section['title']}"
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_LIGHT_GRAY
        p.space_after = Pt(12)


def create_end_slide(prs: Presentation):
    """Create a closing slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_TITLE_BG

    txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "감사합니다"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.5))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Q & A"
    p2.font.size = Pt(24)
    p2.font.color.rgb = COLOR_ACCENT
    p2.alignment = PP_ALIGN.CENTER


def generate_pptx(sections: list[dict], title: str, template_path: str | None, output_path: str):
    """Generate the PPTX file."""
    if template_path and Path(template_path).exists():
        prs = Presentation(template_path)
        print(f"Using template: {template_path}")
    else:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        if template_path:
            print(f"Template not found: {template_path}, using default layout")

    # Title slide
    create_title_slide(
        prs, title,
        subtitle="부품별 기능 · 전압 · 사용법"
    )

    # Table of contents
    titled_sections = [s for s in sections if s["title"]]
    if titled_sections:
        create_toc_slide(prs, titled_sections)

    # Content slides — split sections that are too long
    for section in sections:
        if not section["title"] and not section["content"]:
            continue

        # If section has too many items, split into multiple slides
        content = section["content"]
        MAX_ITEMS_PER_SLIDE = 10

        if len(content) <= MAX_ITEMS_PER_SLIDE:
            create_section_slide(prs, section)
        else:
            # Split into chunks
            for chunk_idx in range(0, len(content), MAX_ITEMS_PER_SLIDE):
                chunk = content[chunk_idx:chunk_idx + MAX_ITEMS_PER_SLIDE]
                suffix = f" ({chunk_idx // MAX_ITEMS_PER_SLIDE + 1})" if len(content) > MAX_ITEMS_PER_SLIDE else ""
                create_section_slide(prs, {
                    "title": section["title"] + suffix,
                    "content": chunk,
                })

    # Closing slide
    create_end_slide(prs)

    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


# ── Main ────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="Notion → PPTX Presentation Generator")
    parser.add_argument("--template", "-t", default=None, help="Path to PPTX template file")
    parser.add_argument("--output", "-o", default="hardware_presentation.pptx", help="Output PPTX path")
    parser.add_argument("--page-id", default=NOTION_PAGE_ID, help="Notion page ID")
    parser.add_argument("--dump-json", action="store_true", help="Dump parsed sections as JSON (debug)")
    args = parser.parse_args()

    if not NOTION_API_KEY:
        print("Error: NOTION_API_KEY environment variable is required.")
        print("Usage: NOTION_API_KEY=ntn_xxx python generate_presentation.py")
        sys.exit(1)

    print("Fetching Notion page content...")
    title = get_page_title(args.page_id)
    print(f"Page title: {title}")

    print("Parsing page blocks...")
    sections = parse_notion_page(args.page_id)
    print(f"Found {len(sections)} sections")

    if args.dump_json:
        print(json.dumps(sections, ensure_ascii=False, indent=2))
        return

    for s in sections:
        items = len(s["content"])
        tables = sum(1 for c in s["content"] if isinstance(c, dict))
        print(f"  - {s['title'] or '(untitled)'}: {items} items ({tables} tables)")

    print("\nGenerating presentation...")
    generate_pptx(sections, title, args.template, args.output)


if __name__ == "__main__":
    main()
